"""
PPTX 생성 서비스

핵심 규칙:
  - 슬라이드 크기 = 원본 이미지와 동일 비율
  - 배경 = 텍스트 제거된 이미지 (슬라이드 전체 채움)
  - 텍스트 박스 = 원본 좌표에서 픽셀→EMU 변환 후 배치
  - 텍스트 박스: fill 없음, 선 없음, 편집 가능
"""

import logging
from pathlib import Path
from typing import List
from lxml import etree

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from app.utils.coordinate_transform import (
    calc_slide_dimensions,
    bbox_to_emu,
    estimate_font_size_pt,
)
from app.utils.color_utils import hex_to_rgb

logger = logging.getLogger(__name__)

# PP_ALIGN 맵
_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def _remove_textbox_border(txBox) -> None:
    """텍스트 박스의 테두리 선을 완전히 제거한다."""
    sp = txBox._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    # 기존 자식 제거 후 noFill 추가
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, qn("a:noFill"))


def _set_textbox_no_fill(txBox) -> None:
    """텍스트 박스 배경 채움 없음."""
    txBox.fill.background()


def _add_text_block(
    slide,
    block: dict,
    img_w: int,
    img_h: int,
    slide_w_emu: int,
    slide_h_emu: int,
) -> None:
    """
    하나의 TextBlock dict를 슬라이드에 editable 텍스트 박스로 추가한다.

    좌표 변환:
      x_emu = bbox.x / img_w * slide_w_emu
      y_emu = bbox.y / img_h * slide_h_emu
      (동일 비율 적용)
    """
    bbox = block.get("bbox", {})
    bx = float(bbox.get("x", 0))
    by = float(bbox.get("y", 0))
    bw = float(bbox.get("w", 10))
    bh = float(bbox.get("h", 10))

    x_emu, y_emu, w_emu, h_emu = bbox_to_emu(bx, by, bw, bh, img_w, img_h, slide_w_emu, slide_h_emu)

    # 텍스트 박스 높이를 폰트 크기 기준으로 최소값 보장
    font_size_pt = float(block.get("fontSize", 16))
    min_h_emu = int(Pt(font_size_pt) * 2)  # 최소 2줄 분량
    h_emu = max(h_emu, min_h_emu)

    # 슬라이드 경계 클램핑
    x_emu = max(0, min(x_emu, slide_w_emu - 1))
    y_emu = max(0, min(y_emu, slide_h_emu - 1))
    w_emu = max(int(Pt(font_size_pt)), min(w_emu, slide_w_emu - x_emu))
    h_emu = min(h_emu, slide_h_emu - y_emu)

    # 텍스트 박스 추가
    txBox = slide.shapes.add_textbox(x_emu, y_emu, w_emu, h_emu)
    _set_textbox_no_fill(txBox)
    _remove_textbox_border(txBox)

    tf = txBox.text_frame
    tf.word_wrap = True

    # 텍스트 라인 분리
    text = block.get("text", "")
    line_breaks: List[str] = block.get("lineBreaks", [])
    lines = line_breaks if (line_breaks and len(line_breaks) >= 1) else [text]
    lines = [l for l in lines if l]  # 빈 문자열 제거
    if not lines:
        lines = [text]

    # 스타일
    bold = bool(block.get("bold", False))
    italic = bool(block.get("italic", False))
    font_family: str = block.get("fontFamily", "Malgun Gothic")
    align_str: str = block.get("align", "left")
    pp_align = _ALIGN_MAP.get(align_str, PP_ALIGN.LEFT)

    color_hex: str = block.get("color", "#000000")
    r, g, b = hex_to_rgb(color_hex)

    # 각 줄을 단락으로 추가
    for i, line_text in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        para.alignment = pp_align

        run = para.add_run()
        run.text = line_text

        font = run.font
        font.name = font_family
        font.size = Pt(font_size_pt)
        font.bold = bold
        font.italic = italic
        font.color.rgb = RGBColor(r, g, b)


def generate_pptx(
    cleaned_image_path: str,
    blocks: List[dict],
    image_width: int,
    image_height: int,
    output_path: str,
) -> None:
    """
    텍스트 제거된 배경 이미지 + editable 텍스트 박스 PPTX를 생성한다.

    Args:
        cleaned_image_path: 텍스트 제거된 배경 이미지 경로
        blocks:             TextBlock dict 리스트 (사용자 편집 반영)
        image_width:        원본 이미지 너비 (px)
        image_height:       원본 이미지 높이 (px)
        output_path:        출력 .pptx 경로
    """
    if not Path(cleaned_image_path).exists():
        raise FileNotFoundError(f"배경 이미지 없음: {cleaned_image_path}")

    prs = Presentation()

    # ── 슬라이드 크기 설정 (원본 이미지 비율 동일) ──────────────────────────
    slide_w_emu, slide_h_emu = calc_slide_dimensions(image_width, image_height)
    prs.slide_width = slide_w_emu
    prs.slide_height = slide_h_emu

    logger.info(
        "슬라이드 크기: %d×%d EMU (이미지: %d×%d px)",
        slide_w_emu, slide_h_emu, image_width, image_height,
    )

    # ── 빈 슬라이드 추가 ─────────────────────────────────────────────────────
    blank_layout = prs.slide_layouts[6]  # 완전히 빈 레이아웃
    slide = prs.slides.add_slide(blank_layout)

    # ── 배경 이미지 추가 (슬라이드 전체) ─────────────────────────────────────
    slide.shapes.add_picture(
        cleaned_image_path,
        left=0,
        top=0,
        width=slide_w_emu,
        height=slide_h_emu,
    )

    # ── 텍스트 박스 추가 ─────────────────────────────────────────────────────
    added = 0
    for block in blocks:
        try:
            _add_text_block(slide, block, image_width, image_height, slide_w_emu, slide_h_emu)
            added += 1
        except Exception as exc:
            logger.warning("블록 [%s] 추가 실패: %s", block.get("id", "?"), exc)

    logger.info("텍스트 박스 %d/%d 추가 완료", added, len(blocks))

    # ── 저장 ─────────────────────────────────────────────────────────────────
    prs.save(output_path)
    logger.info("PPTX 저장 완료: %s", output_path)
